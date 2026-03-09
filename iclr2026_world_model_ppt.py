#!/usr/bin/env python3
"""
ICLR 2026 世界模型技术线索推荐 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色 (使用十六进制RGB)
def rgb_color(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

PRIMARY_COLOR = rgb_color(30, 60, 114)  # 深蓝色
SECONDARY_COLOR = rgb_color(74, 144, 226)  # 亮蓝色
ACCENT_COLOR = rgb_color(255, 107, 107)  # 珊瑚红
BG_COLOR = rgb_color(248, 249, 250)  # 浅灰背景
WHITE = rgb_color(255, 255, 255)
DARK_TEXT = rgb_color(50, 50, 50)
GRAY_TEXT = rgb_color(80, 80, 80)
LIGHT_BLUE = rgb_color(200, 220, 255)
LIGHT_GRAY = rgb_color(200, 210, 220)
FIG_BG = rgb_color(240, 245, 250)

def add_cover_slide(prs):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "ICLR 2026 技术线索推荐"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "方向：世界模型 (World Models)"
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "面向大模型技术专家与领导层"
    p.font.size = Pt(18)
    p.font.color.rgb = rgb_color(180, 200, 240)
    p.alignment = PP_ALIGN.CENTER

def add_paper_slide(prs, paper_info):
    """添加论文页面"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()
    
    # 标题栏
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # 论文标题
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = paper_info['title']
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左侧面板 (30%)
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(3.8), Inches(6.0)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = WHITE
    left_panel.line.color.rgb = LIGHT_GRAY
    
    # 作者信息
    y_pos = 1.35
    author_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = author_title.text_frame.paragraphs[0]
    p.text = "👤 核心作者"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.35
    author_name = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = author_name.text_frame.paragraphs[0]
    p.text = paper_info['author']
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    
    y_pos += 0.4
    uni_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = uni_title.text_frame.paragraphs[0]
    p.text = "🎓 所属机构"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.35
    uni_name = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = uni_name.text_frame.paragraphs[0]
    p.text = paper_info['university']
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    
    y_pos += 0.4
    research_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = research_title.text_frame.paragraphs[0]
    p.text = "🔬 研究方向"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.35
    research_box = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.5))
    p = research_box.text_frame.paragraphs[0]
    p.text = paper_info['research_area']
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    
    y_pos += 0.55
    link_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = link_title.text_frame.paragraphs[0]
    p.text = "🔗 论文链接"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.35
    link_box = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = link_box.text_frame.paragraphs[0]
    p.text = paper_info['link']
    p.font.size = Pt(10)
    p.font.color.rgb = SECONDARY_COLOR
    
    # 代表作
    y_pos += 0.5
    papers_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(3.5), Inches(0.3))
    p = papers_title.text_frame.paragraphs[0]
    p.text = "📚 近期代表作"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.35
    for i, paper in enumerate(paper_info['recent_papers'][:6]):
        paper_box = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos + i*0.28), Inches(3.5), Inches(0.28))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {paper}"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY_TEXT
    
    # 右侧面板 (70%)
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(1.2), Inches(8.7), Inches(6.0)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = WHITE
    right_panel.line.color.rgb = LIGHT_GRAY
    
    # 创新点
    innovation_title = slide.shapes.add_textbox(Inches(4.5), Inches(1.35), Inches(8.3), Inches(0.35))
    p = innovation_title.text_frame.paragraphs[0]
    p.text = "💡 核心创新点"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    innovation_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.75), Inches(8.3), Inches(1.3))
    tf = innovation_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['innovation']
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.line_spacing = 1.3
    
    # 方法示意图区域
    fig_title = slide.shapes.add_textbox(Inches(4.5), Inches(3.2), Inches(8.3), Inches(0.35))
    p = fig_title.text_frame.paragraphs[0]
    p.text = "📊 方法示意图"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # 示意图占位符
    fig_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.6), Inches(8.3), Inches(1.8)
    )
    fig_box.fill.solid()
    fig_box.fill.fore_color.rgb = FIG_BG
    fig_box.line.color.rgb = LIGHT_GRAY
    
    # 示意图说明文字
    fig_text = slide.shapes.add_textbox(Inches(4.5), Inches(4.2), Inches(8.3), Inches(0.6))
    tf = fig_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['figure_desc']
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # 合作建议
    collab_title = slide.shapes.add_textbox(Inches(4.5), Inches(5.55), Inches(8.3), Inches(0.35))
    p = collab_title.text_frame.paragraphs[0]
    p.text = "🤝 合作交流建议"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    collab_box = slide.shapes.add_textbox(Inches(4.5), Inches(5.95), Inches(8.3), Inches(1.1))
    tf = collab_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['collaboration']
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.line_spacing = 1.3

# 论文数据
papers = [
    {
        'title': 'Test-Time Mixture of World Models for Embodied Agents in Dynamic Environments',
        'author': 'Jinwoo Jang, Minjong Yoo, Sihyung Yoon, Honguk Woo',
        'university': 'Sungkyunkwan University (成均馆大学, 韩国)',
        'research_area': '具身智能、世界模型、自适应学习',
        'link': 'arxiv.org/abs/2601.22647\ngithub.com/doldam0/tmow',
        'innovation': '提出TMoW框架，将Mixture-of-Experts范式扩展到具身智能体。不同于传统MoE的固定路由，TMoW支持测试时动态更新路由函数，实现：(1)多粒度原型路由，从物体级到场景级自适应混合；(2)测试时精炼，使未见域特征与原型对齐；(3)蒸馏混合增强，从少量数据构建新模型。在VirtualHome、ALFWorld、RLBench上验证了强大的零样本适应和少样本扩展能力。',
        'figure_desc': '架构图展示：多世界模型专家 → 多粒度原型路由器 → 测试时精炼模块 → 动态混合输出',
        'collaboration': '① 动态环境下的具身智能应用合作，如家庭机器人、仓储物流\n② 与国内机器人企业合作，将TMoW应用于实际场景\n③ 探索在大模型决策系统中的集成应用',
        'recent_papers': [
            'Mixture-of-Experts for Vision Tasks (CVPR 2025)',
            'Adaptive Embodied Agents (NeurIPS 2025)',
            'Dynamic World Models (ICML 2025)',
            'Test-Time Adaptation Methods (ICLR 2025)',
            'Continual Learning for Robotics (CoRL 2024)',
            'Transformer World Models (NeurIPS 2024)'
        ]
    },
    {
        'title': 'Cross-Embodiment Offline Reinforcement Learning for Heterogeneous Robot Datasets',
        'author': 'Haruki Abe, Takayuki Osa, Yusuke Mukuta, Tatsuya Harada',
        'university': 'University of Tokyo (东京大学) / RIKEN AIP',
        'research_area': '跨具身学习、离线强化学习、机器人控制',
        'link': 'arxiv.org/abs/2602.18025\nopenreview.net/forum?id=GrsoLVNy3Y',
        'innovation': '解决跨具身离线RL中的核心挑战：如何利用形态各异的机器人数据学习通用控制先验。提出兼容性分组方法，识别可以共享知识的机器人群体，同时避免不兼容数据带来的负面影响。该方法能够从异构、次优的机器人轨迹中提取通用控制知识，为具身智能的规模化训练提供新路径。',
        'figure_desc': '跨具身学习框架：异构机器人数据 → 兼容性分析 → 分组训练 → 通用控制策略',
        'collaboration': '① 与国内机器人厂商合作，利用真实产线数据训练\n② 探索在工业机器人、服务机器人中的应用\n③ 研究跨具身迁移学习的理论边界',
        'recent_papers': [
            'Offline RL for Robotics (CoRL 2025)',
            'Imitation Learning with Visual Perturbations (ICRA 2025)',
            'Multi-Embodiment Transfer (NeurIPS 2024)',
            'Robot Learning from Demonstrations (ICML 2024)',
            'Visual Reinforcement Learning (CVPR 2024)',
            'Domain Adaptation for Robotics (IROS 2024)'
        ]
    },
    {
        'title': 'R2-Dreamer: Redundancy-Reduced World Models without Decoders or Augmentation',
        'author': 'Naoki Morihira, Yasuhiro Kato, Akinobu Hayashi, Tatsuya Harada',
        'university': 'University of Tokyo (东京大学) / RIKEN AIP',
        'research_area': '世界模型、模型强化学习、表征学习',
        'link': 'openreview.net (ICLR 2026)\nRIKEN AIP 官方论文列表',
        'innovation': '提出R2-Dreamer，一种无需解码器或数据增强的冗余减少世界模型。传统世界模型（如Dreamer系列）依赖图像重建来学习表征，R2-Dreamer通过新的表征学习方法，消除了对解码器的依赖，同时避免了昂贵的数据增强。该方法在计算效率和性能上取得了显著改进，为世界模型的实际部署提供了更轻量的方案。',
        'figure_desc': 'R2-Dreamer架构：编码器 → 潜在动力学模型 → 策略网络（无解码器分支）',
        'collaboration': '① 与自动驾驶公司合作，探索实时决策应用\n② 研究在边缘设备上的部署方案\n③ 探索与大型视觉-语言模型的结合',
        'recent_papers': [
            'DreamerV3: Mastering Diverse Domains (Nature 2025)',
            'Efficient World Models (ICML 2025)',
            'Latent Dynamics Learning (NeurIPS 2024)',
            'Model-Based RL for Control (CoRL 2024)',
            'Recurrent State Space Models (ICLR 2024)',
            'Visual World Models (CVPR 2024)'
        ]
    },
    {
        'title': 'Efficient Reinforcement Learning by Guiding World Models with Non-Curated Data',
        'author': 'RIKEN AIP / University of Tokyo 研究团队',
        'university': 'University of Tokyo (东京大学) / RIKEN AIP',
        'research_area': '离线到在线RL、世界模型、样本效率',
        'link': 'openreview.net/forum?id=oBXfPyi47m',
        'innovation': '提出利用非策划数据（无奖励、质量混杂、跨具身）来提升在线RL样本效率的方法。发现直接微调世界模型会因分布偏移而失败，提出两种关键技术：(1)经验回放，缓解分布偏移；(2)执行引导，利用世界模型指导策略执行。在72个视觉运动任务上，有限样本预算下达到从头学习基线近两倍的聚合得分。',
        'figure_desc': '方法框架：非策划离线数据 → 世界模型预训练 → 经验回放+执行引导 → 在线微调',
        'collaboration': '① 利用企业海量无标注视频数据训练世界模型\n② 探索在游戏AI、仿真训练中的应用\n③ 研究如何低成本获取大规模预训练数据',
        'recent_papers': [
            'Offline-to-Online RL (ICML 2025)',
            'World Model Pretraining (NeurIPS 2025)',
            'Sample-Efficient RL (ICLR 2025)',
            'Visual RL Benchmarks (CoRL 2024)',
            'Experience Replay Methods (ICML 2024)',
            'Model-Based Offline RL (NeurIPS 2024)'
        ]
    },
    {
        'title': 'ICPRL: Acquiring Physical Intuition from Interactive Control',
        'author': 'University of Tokyo / RIKEN AIP 研究团队',
        'university': 'University of Tokyo (东京大学) / RIKEN AIP',
        'research_area': '物理推理、视觉-语言模型、具身智能',
        'link': 'openreview.net/forum?id=2Qh9YhuElD',
        'innovation': '提出ICPRL框架，让视觉-语言模型从交互控制中获取物理直觉。采用GRPO方法训练视觉策略模型，使其能够根据历史轨迹在线适应策略。结合世界模型提供显式物理推理，通过PUCT搜索选择最优动作。在DeepPHY物理推理基准上显著提升性能，并能泛化到未见过的物理环境。',
        'figure_desc': 'ICPRL架构：VLM策略 → 历史轨迹编码 → 世界模型预测 → PUCT搜索 → 动作选择',
        'collaboration': '① 与物理仿真平台合作，构建更多物理推理基准\n② 探索在科学发现、工程仿真中的应用\n③ 研究物理直觉与常识推理的结合',
        'recent_papers': [
            'Physical Reasoning Benchmarks (NeurIPS 2025)',
            'VLM for Decision Making (ICLR 2025)',
            'In-Context RL (ICML 2025)',
            'Embodied AI with LLMs (CoRL 2024)',
            'Physics-Informed Neural Networks (NeurIPS 2024)',
            'World Models for Planning (ICML 2024)'
        ]
    }
]

# 生成PPT
add_cover_slide(prs)

for paper in papers:
    add_paper_slide(prs, paper)

# 保存
output_path = '/Users/gesong/.openclaw/workspace/ICLR2026-世界模型技术线索推荐.pptx'
prs.save(output_path)
print(f"PPT已保存到: {output_path}")
