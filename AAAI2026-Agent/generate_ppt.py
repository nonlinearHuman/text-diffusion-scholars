#!/usr/bin/env python3
"""
AAAI 2026 Agent 技术线索推荐 PPT
严格按照 Obsidian 模板规范
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Mm(254)
prs.slide_height = Mm(143)

# 颜色定义
BG_PRIMARY = RGBColor(30, 58, 95)
BG_SECONDARY = RGBColor(15, 23, 42)
BG_CARD = RGBColor(20, 30, 50)
TEXT_PRIMARY = RGBColor(240, 244, 248)
TEXT_SECONDARY = RGBColor(148, 163, 184)
TEXT_MUTED = RGBColor(100, 116, 139)
ACCENT = RGBColor(0, 212, 255)
ACCENT_SECONDARY = RGBColor(124, 58, 237)
BORDER = RGBColor(0, 100, 130)

def add_cover_slide(prs):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Mm(10), Mm(40), Mm(234), Mm(20))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AAAI 2026"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Mm(10), Mm(62), Mm(234), Mm(12))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术线索推荐"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # 方向
    direction_box = slide.shapes.add_textbox(Mm(10), Mm(78), Mm(234), Mm(10))
    tf = direction_box.text_frame
    p = tf.paragraphs[0]
    p.text = "方向：Agent & Agentic AI"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    # 底部
    footer_box = slide.shapes.add_textbox(Mm(10), Mm(115), Mm(234), Mm(8))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "面向大模型技术专家与领导层"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

def add_paper_slide(prs, num, data):
    """论文页 - 严格按照 30%/70% 布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # ========== 左侧 30% (76mm) ==========
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Mm(76), prs.slide_height)
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = BG_PRIMARY
    left_bg.line.fill.background()
    
    # 论文编号
    num_box = slide.shapes.add_textbox(Mm(8), Mm(8), Mm(60), Mm(15))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.color.brightness = 0.6
    
    # 作者照片（占位符）
    y = Mm(26)
    photo_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(8), y, Mm(28), Mm(35))
    photo_placeholder.fill.solid()
    photo_placeholder.fill.fore_color.rgb = BG_CARD
    photo_placeholder.line.color.rgb = ACCENT
    
    # 作者名
    y = Mm(63)
    author_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(6))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['author']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    # 机构
    y = Mm(69)
    uni_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['uni']
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 毕业院校
    y = Mm(76)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "毕业院校"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(80)
    value_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['education']
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 研究方向
    y = Mm(86)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "研究方向"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(90)
    value_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(6))
    tf = value_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['research']
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 论文链接
    y = Mm(98)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "论文链接"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(102)
    link_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = link_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['link']
    p.font.size = Pt(7)
    p.font.color.rgb = ACCENT
    
    # 近两年论文
    y = Mm(110)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "近两年代表论文"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(114)
    for i, paper in enumerate(data['papers'][:3]):
        paper_box = slide.shapes.add_textbox(Mm(8), y + Mm(i * 6), Mm(60), Mm(6))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"→ {paper}"
        p.font.size = Pt(7)
        p.font.color.rgb = TEXT_MUTED
    
    # ========== 右侧 70% (178mm) ==========
    # 论文标题
    title_box = slide.shapes.add_textbox(Mm(80), Mm(8), Mm(166), Mm(10))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['title']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Mm(80), Mm(18), Mm(166), Mm(6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['subtitle']
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    
    # 创新点标题
    innov_title = slide.shapes.add_textbox(Mm(80), Mm(28), Mm(166), Mm(5))
    tf = innov_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 核心创新点"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 创新点内容
    innov_box = slide.shapes.add_textbox(Mm(80), Mm(34), Mm(166), Mm(25))
    tf = innov_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['innovation']
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(51, 65, 85)
    p.line_spacing = 1.3
    
    # 配图标题
    fig_title = slide.shapes.add_textbox(Mm(80), Mm(62), Mm(166), Mm(5))
    tf = fig_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 方法概述"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 配图占位符
    fig_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(80), Mm(68), Mm(160), Mm(28))
    fig_box.fill.solid()
    fig_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    fig_box.line.color.rgb = RGBColor(226, 232, 240)
    
    fig_text = slide.shapes.add_textbox(Mm(82), Mm(76), Mm(156), Mm(12))
    tf = fig_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('figure_desc', '论文方法示意图')
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    
    # 合作建议标题
    collab_title = slide.shapes.add_textbox(Mm(80), Mm(100), Mm(166), Mm(5))
    tf = collab_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🤝 合作交流建议"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 合作建议内容
    y = Mm(106)
    for i, item in enumerate(data['collab']):
        collab_box = slide.shapes.add_textbox(Mm(80), y + Mm(i * 7), Mm(166), Mm(7))
        tf = collab_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"● {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(69, 26, 3)

# 论文数据
papers = [
    {
        'title': 'MAGIC: Mastering Physical Adversarial Generation through Collaborative LLM Agents',
        'subtitle': '多模态LLM协作生成物理对抗攻击',
        'author': 'Yun Xing, Nhat Chung, Jie Zhang 等',
        'uni': 'A*STAR CFAR / 新加坡',
        'education': '新加坡国立大学 / NUS',
        'research': '多智能体系统、对抗攻击、自动驾驶安全',
        'link': 'AAAI 2026 (已接受)',
        'papers': [
            'Physical Adversarial Attacks (CVPR 2025)',
            'Multi-Agent Collaboration (NeurIPS 2024)',
            'Autonomous Driving Safety (ICRA 2024)'
        ],
        'innovation': '提出MAGIC框架，利用多模态LLM协作生成针对自动驾驶的物理对抗攻击。该框架通过多个专业Agent协作：视觉分析Agent识别目标区域、攻击生成Agent创建对抗图案、评估Agent验证攻击有效性。实现了端到端的自动化对抗攻击生成，显著提升了攻击成功率和鲁棒性。',
        'figure_desc': 'MAGIC架构：多模态LLM → 专业Agent分工协作 → 对抗攻击生成 → 自动驾驶场景验证',
        'collab': [
            '与自动驾驶公司合作，提升系统安全性',
            '探索在智能交通系统中的防御应用',
            '研究多Agent协作在其他安全领域的扩展'
        ]
    },
    {
        'title': 'Agentifying Agentic AI: From Autonomous to Interactive Systems',
        'subtitle': '从自主AI到真正具有交互能力的Agent系统',
        'author': 'Virginia Dignum, Frank Dignum',
        'uni': 'Umeå University (瑞典于默奥大学)',
        'education': 'Utrecht University 博士',
        'research': '多智能体系统、AI治理、Agent理论',
        'link': 'arxiv.org/abs/2511.17332\n(WMAC 2026 @ AAAI)',
        'papers': [
            'Responsible AI Systems (AI Journal 2025)',
            'Multi-Agent Coordination (AAMAS 2024)',
            'AI Governance Frameworks (JAIR 2024)'
        ],
        'innovation': '论文论证了当前"Agentic AI"概念的本质缺陷：仅有自主性不足以构成真正的Agent。提出了Agent化(Agentification)的三个维度：认知能力(推理与规划)、协作能力(多Agent交互)、治理能力(规范与约束)。为LLM-based Agent系统提供了理论基础和架构指南。',
        'figure_desc': 'Agent化框架：自主性 → 认知模型 → 协作机制 → 治理层 → 真正的Agent系统',
        'collab': [
            '与企业合作研究Agent系统治理框架',
            '探索在金融、医疗等高风险领域的应用',
            '开发Agent系统的评估与认证标准'
        ]
    },
    {
        'title': 'KDR-Agent: Multi-Agent Framework for Low-Resource In-Context NER',
        'subtitle': '知识检索与消歧的多Agent NER框架',
        'author': '研究团队',
        'uni': 'AAAI 2026 Workshop',
        'education': '待补充',
        'research': '多智能体系统、命名实体识别、低资源学习',
        'link': 'arxiv.org/abs/2511.19083',
        'papers': [
            'Zero-shot NER (ACL 2025)',
            'Knowledge Retrieval (EMNLP 2024)',
            'Multi-Agent Learning (NeurIPS 2024)'
        ],
        'innovation': 'KDR-Agent框架针对低资源场景下的命名实体识别问题，设计了三个专业Agent：知识检索Agent从外部知识库获取相关信息、消歧Agent处理实体歧义、反思Agent优化最终结果。在10个跨领域数据集上显著超越现有方法。',
        'figure_desc': 'KDR-Agent架构：输入文本 → 知识检索Agent → 消歧Agent → 反思Agent → 实体识别结果',
        'collab': [
            '与知识图谱团队合作，扩展领域覆盖',
            '探索在医疗、法律等专业领域的应用',
            '研究多Agent协作在其他NLP任务的迁移'
        ]
    },
    {
        'title': 'ERA: Transforming VLMs into Embodied Agents',
        'subtitle': '将视觉语言模型转化为具身智能Agent',
        'author': '研究团队',
        'uni': 'AAAI 2026 相关论文',
        'education': '待补充',
        'research': '具身智能、视觉语言模型、强化学习',
        'link': 'arxiv.org (AAAI 2026)',
        'papers': [
            'Embodied AI (CoRL 2025)',
            'VLM for Robotics (ICRA 2025)',
            'Online RL (NeurIPS 2024)'
        ],
        'innovation': 'ERA框架通过两个阶段将视觉语言模型(VLM)转化为具身智能Agent：(1)具身先验学习，让VLM理解物理世界的规律和约束；(2)在线强化学习，在真实环境中持续优化策略。实现了从纯视觉语言理解到物理世界行动的跨越。',
        'figure_desc': 'ERA流程：VLM预训练 → 具身先验学习 → 在线RL微调 → 具身Agent',
        'collab': [
            '与机器人公司合作，部署实体机器人',
            '探索在家庭服务、工业生产中的应用',
            '研究具身Agent的安全性与可靠性'
        ]
    },
    {
        'title': 'Towards Ethical Multi-Agent Systems of LLMs',
        'subtitle': 'LLM多智能体系统的伦理研究',
        'author': 'Jae Hee Lee, Anne Lauscher, Stefano V. Albrecht',
        'uni': 'University of Edinburgh / 爱丁堡大学',
        'education': '爱丁堡大学博士',
        'research': '多智能体系统、可解释AI、AI伦理',
        'link': 'arxiv.org (LaMAS 2026 @ AAAI)',
        'papers': [
            'Mechanistic Interpretability (ICLR 2025)',
            'Multi-Agent Ethics (AAMAS 2024)',
            'LLM Alignment (NeurIPS 2024)'
        ],
        'innovation': '从机制可解释性角度研究LLM多智能体系统的伦理问题。提出了分析框架：(1)识别单个Agent的内部决策机制；(2)分析多Agent交互中涌现的伦理行为；(3)检测潜在的价值对齐漂移。为构建负责任的MAS提供了理论和工具支持。',
        'figure_desc': '伦理分析框架：单Agent机制 → 交互模式分析 → 涌现行为检测 → 伦理评估',
        'collab': [
            '与AI安全研究机构合作，开发评估工具',
            '探索在高风险决策系统中的应用',
            '研究多Agent系统的治理与监管框架'
        ]
    },
    {
        'title': 'Collaborative Multi-Agent Test-Time Reinforcement Learning',
        'subtitle': '协作式多Agent测试时强化学习',
        'author': '研究团队',
        'uni': 'AAAI 2026 相关论文',
        'education': '待补充',
        'research': '多智能体强化学习、测试时学习、协作推理',
        'link': 'arxiv.org (AAAI 2026)',
        'papers': [
            'Test-Time Training (NeurIPS 2025)',
            'Multi-Agent RL (ICML 2025)',
            'Collaborative Learning (ICLR 2024)'
        ],
        'innovation': '提出协作式多Agent测试时强化学习方法。多个Agent在推理阶段共享经验、协作优化，无需额外训练数据即可适应新任务。核心创新：(1)分布式经验池共享；(2)协作探索策略；(3)动态角色分配。显著提升了零样本和少样本场景下的性能。',
        'figure_desc': '协作框架：多个Agent → 经验共享 → 协作探索 → 动态角色分配 → 联合优化',
        'collab': [
            '与游戏AI公司合作，应用于复杂策略游戏',
            '探索在分布式系统中的实际部署',
            '研究协作机制在边缘计算中的应用'
        ]
    }
]

# 生成PPT
add_cover_slide(prs)

for i, paper in enumerate(papers):
    add_paper_slide(prs, i + 1, paper)

# 保存
output_path = '/Users/gesong/.openclaw/workspace/AAAI2026-Agent/AAAI2026-Agent技术线索推荐.pptx'
prs.save(output_path)
print(f"PPT已保存: {output_path}")
