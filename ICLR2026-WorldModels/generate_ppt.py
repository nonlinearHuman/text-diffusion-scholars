#!/usr/bin/env python3
"""
ICLR 2026 世界模型技术线索推荐 PPT
严格按照 Obsidian 模板规范
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Mm(254)
prs.slide_height = Mm(143)

# 颜色定义
BG_PRIMARY = RGBColor(30, 58, 95)
BG_SECONDARY = RGBColor(15, 23, 42)
BG_CARD = RGBColor(20, 30, 50)
TEXT_PRIMARY = RGBColor(240, 244, 248)
TEXT_SECONDARY = RGBColor(148, 163, 184)
TEXT_MUTED = RGBColor(100, 116, 139)
ACCENT = RGBColor(0, 212, 255)
ACCENT_SECONDARY = RGBColor(124, 58, 237)
BORDER = RGBColor(0, 100, 130)

def add_cover_slide(prs):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Mm(10), Mm(40), Mm(234), Mm(20))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ICLR 2026"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Mm(10), Mm(62), Mm(234), Mm(12))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术线索推荐"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # 方向
    direction_box = slide.shapes.add_textbox(Mm(10), Mm(78), Mm(234), Mm(10))
    tf = direction_box.text_frame
    p = tf.paragraphs[0]
    p.text = "方向：世界模型 (World Models)"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    # 底部
    footer_box = slide.shapes.add_textbox(Mm(10), Mm(115), Mm(234), Mm(8))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "面向大模型技术专家与领导层"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

def add_paper_slide(prs, num, data):
    """论文页 - 严格按照 30%/70% 布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # ========== 左侧 30% (76mm) ==========
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Mm(76), prs.slide_height)
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = BG_PRIMARY
    left_bg.line.fill.background()
    
    # 论文编号
    num_box = slide.shapes.add_textbox(Mm(8), Mm(8), Mm(60), Mm(15))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.color.brightness = 0.6
    
    # 作者照片
    y = Mm(26)
    if data.get('photo') and os.path.exists(data['photo']):
        try:
            photo = slide.shapes.add_picture(data['photo'], Mm(8), y, width=Mm(28), height=Mm(35))
        except:
            # 照片加载失败时用占位符
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(8), y, Mm(28), Mm(35))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = BG_CARD
            placeholder.line.color.rgb = ACCENT
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(8), y, Mm(28), Mm(35))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_CARD
        placeholder.line.color.rgb = ACCENT
    
    # 作者名
    y = Mm(63)
    author_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(6))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['author']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    # 机构
    y = Mm(69)
    uni_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['uni']
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 毕业院校
    y = Mm(76)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "毕业院校"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(80)
    value_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['education']
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 研究方向
    y = Mm(86)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "研究方向"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(90)
    value_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(6))
    tf = value_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['research']
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 论文链接
    y = Mm(98)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "论文链接"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(102)
    link_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = link_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['link']
    p.font.size = Pt(7)
    p.font.color.rgb = ACCENT
    
    # 近两年论文
    y = Mm(110)
    label_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "近两年代表论文"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    y = Mm(114)
    for i, paper in enumerate(data['papers'][:3]):
        paper_box = slide.shapes.add_textbox(Mm(8), y + Mm(i * 6), Mm(60), Mm(6))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"→ {paper}"
        p.font.size = Pt(7)
        p.font.color.rgb = TEXT_MUTED
    
    # ========== 右侧 70% (178mm) ==========
    # 论文标题
    title_box = slide.shapes.add_textbox(Mm(80), Mm(8), Mm(166), Mm(10))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['title']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Mm(80), Mm(18), Mm(166), Mm(6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['subtitle']
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    
    # 创新点标题
    innov_title = slide.shapes.add_textbox(Mm(80), Mm(28), Mm(166), Mm(5))
    tf = innov_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 核心创新点"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 创新点内容
    innov_box = slide.shapes.add_textbox(Mm(80), Mm(34), Mm(166), Mm(20))
    tf = innov_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['innovation']
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(51, 65, 85)
    
    # 配图标题
    fig_title = slide.shapes.add_textbox(Mm(80), Mm(56), Mm(166), Mm(5))
    tf = fig_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 论文配图"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 配图
    if data.get('figure') and os.path.exists(data['figure']):
        try:
            fig = slide.shapes.add_picture(data['figure'], Mm(80), Mm(62), height=Mm(35))
            # 如果图片太宽，调整
            if fig.width > Mm(160):
                fig.width = Mm(160)
        except:
            pass
    
    # 合作建议标题
    collab_title = slide.shapes.add_textbox(Mm(80), Mm(100), Mm(166), Mm(5))
    tf = collab_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🤝 合作交流建议"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # 合作建议内容
    y = Mm(106)
    for i, item in enumerate(data['collab']):
        collab_box = slide.shapes.add_textbox(Mm(80), y + Mm(i * 7), Mm(166), Mm(7))
        tf = collab_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"● {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(69, 26, 3)

# 论文数据
papers = [
    {
        'title': 'Test-Time Mixture of World Models',
        'subtitle': 'TMoW: 动态环境下的具身智能自适应框架',
        'author': 'Honguk Woo 等',
        'uni': '成均馆大学 (SKKU)',
        'education': 'UT Austin 博士',
        'research': '具身智能、世界模型、自适应学习',
        'link': 'arxiv.org/abs/2601.22647',
        'papers': [
            'Incremental Learning of Skills (NeurIPS 2024)',
            'Embodied CoT Distillation (2024)',
            'Global DAG Task Scheduler (2021)'
        ],
        'innovation': '提出TMoW框架，将MoE扩展到具身智能体。支持测试时动态更新路由函数，实现多粒度原型路由、测试时精炼和蒸馏混合增强。在VirtualHome、ALFWorld、RLBench上验证了强大的零样本适应能力，比SayCanPay提升27.21%。',
        'collab': [
            '动态环境下的具身智能应用合作（家庭机器人、仓储物流）',
            '与国内机器人企业合作，将TMoW应用于实际场景',
            '探索在大模型决策系统中的集成应用'
        ],
        'photo': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/woo_photo.jpg',
        'figure': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/tmow_fig2.png'
    },
    {
        'title': 'Cross-Embodiment Offline RL for Heterogeneous Robot Datasets',
        'subtitle': '异构机器人数据集的跨具身离线强化学习',
        'author': 'Haruki Abe, Tatsuya Harada 等',
        'uni': '东京大学 / RIKEN AIP',
        'education': '东京大学博士 (2001)',
        'research': '跨具身学习、离线强化学习、机器人控制',
        'link': 'arxiv.org/abs/2602.18025',
        'papers': [
            'Offline RL for Robotics (CoRL 2025)',
            'Multi-Embodiment Transfer (NeurIPS 2024)',
            'Visual RL (CVPR 2024)'
        ],
        'innovation': '解决跨具身离线RL中如何利用形态各异的机器人数据学习通用控制先验的问题。提出基于形态相似性的分组策略，减少机器人间的梯度冲突，从异构次优轨迹中提取通用控制知识，为具身智能规模化训练提供新路径。',
        'collab': [
            '与国内机器人厂商合作，利用真实产线数据训练',
            '探索在工业机器人、服务机器人中的应用',
            '研究跨具身迁移学习的理论边界'
        ],
        'photo': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/harada_photo.jpg',
        'figure': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/cross_embodiment_fig5.png'
    },
    {
        'title': 'R2-Dreamer: Redundancy-Reduced World Models',
        'subtitle': '无解码器的冗余减少世界模型',
        'author': 'Naoki Morihira, Tatsuya Harada 等',
        'uni': '东京大学 / RIKEN AIP',
        'education': '东京大学',
        'research': '世界模型、模型强化学习、表征学习',
        'link': 'openreview.net (ICLR 2026)',
        'papers': [
            'DreamerV3 (Nature 2025)',
            'Efficient World Models (ICML 2025)',
            'Latent Dynamics (NeurIPS 2024)'
        ],
        'innovation': '提出无需解码器或数据增强的冗余减少世界模型。消除了对解码器的依赖，避免昂贵的数据增强，在计算效率和性能上取得显著改进，为世界模型的实际部署提供更轻量方案。',
        'collab': [
            '与自动驾驶公司合作，探索实时决策应用',
            '研究在边缘设备上的部署方案',
            '探索与大型视觉-语言模型的结合'
        ],
        'photo': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/harada_photo.jpg',
        'figure': None
    },
    {
        'title': 'Efficient RL by Guiding World Models with Non-Curated Data',
        'subtitle': '利用非策划数据提升在线RL样本效率',
        'author': 'RIKEN AIP / U-Tokyo 团队',
        'uni': '东京大学 / RIKEN AIP',
        'education': '东京大学 / RIKEN',
        'research': '离线到在线RL、世界模型、样本效率',
        'link': 'openreview.net/forum?id=oBXfPyi47m',
        'papers': [
            'Offline-to-Online RL (ICML 2025)',
            'World Model Pretraining (NeurIPS 2025)',
            'Sample-Efficient RL (ICLR 2025)'
        ],
        'innovation': '利用非策划数据（无奖励、质量混杂、跨具身）提升在线RL样本效率。提出经验回放和执行引导两种技术，在72个视觉运动任务上达到从头学习基线近两倍的聚合得分。',
        'collab': [
            '利用企业海量无标注视频数据训练世界模型',
            '探索在游戏AI、仿真训练中的应用',
            '研究如何低成本获取大规模预训练数据'
        ],
        'photo': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/harada_photo.jpg',
        'figure': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/guiding_fig.jpeg'
    },
    {
        'title': 'ICPRL: Acquiring Physical Intuition from Interactive Control',
        'subtitle': '从交互控制中获取物理直觉',
        'author': 'RIKEN AIP / U-Tokyo 团队',
        'uni': '东京大学 / RIKEN AIP',
        'education': '东京大学 / RIKEN',
        'research': '物理推理、视觉-语言模型、具身智能',
        'link': 'openreview.net/forum?id=2Qh9YhuElD',
        'papers': [
            'Physical Reasoning Benchmarks (NeurIPS 2025)',
            'VLM for Decision Making (ICLR 2025)',
            'In-Context RL (ICML 2025)'
        ],
        'innovation': '让视觉-语言模型从交互控制中获取物理直觉。采用GRPO方法训练视觉策略模型，结合世界模型提供显式物理推理，通过PUCT搜索选择最优动作，显著提升物理推理能力。',
        'collab': [
            '与物理仿真平台合作，构建更多物理推理基准',
            '探索在科学发现、工程仿真中的应用',
            '研究物理直觉与常识推理的结合'
        ],
        'photo': '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/paper_figures/harada_photo.jpg',
        'figure': None
    }
]

# 生成PPT
add_cover_slide(prs)

for i, paper in enumerate(papers):
    add_paper_slide(prs, i + 1, paper)

# 保存
output_path = '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/ICLR2026-世界模型技术线索推荐.pptx'
prs.save(output_path)
print(f"PPT已保存: {output_path}")
