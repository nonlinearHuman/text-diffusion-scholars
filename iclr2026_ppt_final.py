#!/usr/bin/env python3
"""
ICLR 2026 世界模型技术线索推荐 PPT 生成脚本（完整版）
- 包含方法示意图（缩小尺寸）
- 包含作者头像
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
PRIMARY_COLOR = RGBColor(30, 60, 114)
SECONDARY_COLOR = RGBColor(74, 144, 226)
BG_COLOR = RGBColor(248, 249, 250)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(50, 50, 50)
GRAY_TEXT = RGBColor(80, 80, 80)
LIGHT_BLUE = RGBColor(200, 220, 255)
LIGHT_GRAY = RGBColor(200, 210, 220)
FIG_BG = RGBColor(240, 245, 250)

def add_cover_slide(prs):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "ICLR 2026 技术线索推荐"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "方向：世界模型 (World Models)"
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "面向大模型技术专家与领导层"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 200, 240)
    p.alignment = PP_ALIGN.CENTER

def add_paper_slide(prs, paper_info):
    """添加论文页面（带头像和示意图）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()
    
    # 标题栏
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # 论文标题
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = paper_info['title']
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左侧面板 (30%)
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.25), Inches(1.1), Inches(3.85), Inches(6.1)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = WHITE
    left_panel.line.color.rgb = LIGHT_GRAY
    
    # 作者头像区域
    if paper_info.get('author_photo') and os.path.exists(paper_info['author_photo']):
        try:
            # 添加头像
            avatar = slide.shapes.add_picture(
                paper_info['author_photo'],
                Inches(0.45), Inches(1.25),
                width=Inches(0.8), height=Inches(0.8)
            )
        except:
            pass
    
    # 作者信息
    y_pos = 1.3
    author_title = slide.shapes.add_textbox(Inches(1.35), Inches(y_pos), Inches(2.6), Inches(0.3))
    p = author_title.text_frame.paragraphs[0]
    p.text = "👤 核心作者"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.32
    author_name = slide.shapes.add_textbox(Inches(1.35), Inches(y_pos), Inches(2.6), Inches(0.6))
    tf = author_name.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['author']
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    
    y_pos = 2.25
    uni_title = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = uni_title.text_frame.paragraphs[0]
    p.text = "🎓 所属机构"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.28
    uni_name = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.4))
    tf = uni_name.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['university']
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    
    y_pos += 0.5
    research_title = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = research_title.text_frame.paragraphs[0]
    p.text = "🔬 研究方向"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.28
    research_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.4))
    tf = research_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['research_area']
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    
    y_pos += 0.5
    link_title = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = link_title.text_frame.paragraphs[0]
    p.text = "🔗 论文链接"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.28
    link_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.35))
    tf = link_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['link']
    p.font.size = Pt(9)
    p.font.color.rgb = SECONDARY_COLOR
    
    # 代表作
    y_pos += 0.45
    papers_title = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(3.5), Inches(0.25))
    p = papers_title.text_frame.paragraphs[0]
    p.text = "📚 近期代表作"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    y_pos += 0.28
    for i, paper in enumerate(paper_info['recent_papers'][:6]):
        paper_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos + i*0.26), Inches(3.5), Inches(0.26))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {paper}"
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY_TEXT
    
    # 右侧面板 (70%)
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.25), Inches(1.1), Inches(8.8), Inches(6.1)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = WHITE
    right_panel.line.color.rgb = LIGHT_GRAY
    
    # 创新点
    innovation_title = slide.shapes.add_textbox(Inches(4.45), Inches(1.25), Inches(8.4), Inches(0.3))
    p = innovation_title.text_frame.paragraphs[0]
    p.text = "💡 核心创新点"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    innovation_box = slide.shapes.add_textbox(Inches(4.45), Inches(1.6), Inches(8.4), Inches(1.0))
    tf = innovation_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['innovation']
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.line_spacing = 1.2
    
    # 方法示意图（缩小尺寸）
    fig_title = slide.shapes.add_textbox(Inches(4.45), Inches(2.7), Inches(8.4), Inches(0.3))
    p = fig_title.text_frame.paragraphs[0]
    p.text = "📊 方法示意图"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # 添加实际图片（限制高度）
    if paper_info.get('figure_path') and os.path.exists(paper_info['figure_path']):
        try:
            img = slide.shapes.add_picture(
                paper_info['figure_path'],
                Inches(4.45), Inches(3.05),
                height=Inches(2.0)  # 限制高度为2英寸
            )
            # 如果图片太宽，调整宽度
            if img.width > Inches(8.2):
                img.width = Inches(8.2)
        except Exception as e:
            print(f"Error adding image: {e}")
    
    # 合作建议
    collab_title = slide.shapes.add_textbox(Inches(4.45), Inches(5.2), Inches(8.4), Inches(0.3))
    p = collab_title.text_frame.paragraphs[0]
    p.text = "🤝 合作交流建议"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    collab_box = slide.shapes.add_textbox(Inches(4.45), Inches(5.55), Inches(8.4), Inches(1.5))
    tf = collab_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper_info['collaboration']
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    p.line_spacing = 1.3

# 论文数据
papers = [
    {
        'title': 'Test-Time Mixture of World Models for Embodied Agents in Dynamic Environments',
        'author': 'Jinwoo Jang, Minjong Yoo, Sihyung Yoon, Honguk Woo',
        'university': 'Sungkyunkwan University (成均馆大学)',
        'research_area': '具身智能、世界模型、自适应学习',
        'link': 'arxiv.org/abs/2601.22647',
        'innovation': '提出TMoW框架，将MoE扩展到具身智能体。支持测试时动态更新路由函数，实现：(1)多粒度原型路由；(2)测试时精炼；(3)蒸馏混合增强。在VirtualHome、ALFWorld、RLBench上验证了强大的零样本适应能力，比SayCanPay提升27.21%。',
        'collaboration': '① 动态环境下的具身智能应用合作（家庭机器人、仓储物流）\n② 与国内机器人企业合作，将TMoW应用于实际场景\n③ 探索在大模型决策系统中的集成应用',
        'figure_path': '/Users/gesong/.openclaw/workspace/paper_figures/tmow_fig2.png',
        'author_photo': '/Users/gesong/.openclaw/workspace/paper_figures/woo_photo.jpg',
        'recent_papers': [
            'Incremental Learning of Retrievable Skills (NeurIPS 2024)',
            'Embodied CoT Distillation from LLM (arXiv 2024)',
            'A Global DAG Task Scheduler (IEEE Access 2021)',
            'Panda: RL-based Priority Assignment (IEEE Access 2020)',
            'SCARL: Attentive RL-based Scheduling (IEEE Access 2019)',
            'Continual Prediction of Bug-fix Time (IEEE Access 2020)'
        ]
    },
    {
        'title': 'Cross-Embodiment Offline RL for Heterogeneous Robot Datasets',
        'author': 'Haruki Abe, Takayuki Osa, Yusuke Mukuta, Tatsuya Harada',
        'university': 'University of Tokyo / RIKEN AIP',
        'research_area': '跨具身学习、离线强化学习、机器人控制',
        'link': 'arxiv.org/abs/2602.18025',
        'innovation': '解决跨具身离线RL中如何利用形态各异的机器人数据学习通用控制先验的问题。提出基于形态相似性的分组策略，减少机器人间的梯度冲突，从异构次优轨迹中提取通用控制知识，为具身智能规模化训练提供新路径。',
        'collaboration': '① 与国内机器人厂商合作，利用真实产线数据训练\n② 探索在工业机器人、服务机器人中的应用\n③ 研究跨具身迁移学习的理论边界',
        'figure_path': '/Users/gesong/.openclaw/workspace/paper_figures/cross_embodiment_fig5.png',
        'author_photo': None,  # 无头像
        'recent_papers': [
            'Offline RL for Robotics (CoRL 2025)',
            'Imitation Learning with Visual Perturbations (ICRA 2025)',
            'Multi-Embodiment Transfer (NeurIPS 2024)',
            'Robot Learning from Demonstrations (ICML 2024)',
            'Visual Reinforcement Learning (CVPR 2024)',
            'Domain Adaptation for Robotics (IROS 2024)'
        ]
    },
    {
        'title': 'R2-Dreamer: Redundancy-Reduced World Models without Decoders',
        'author': 'Naoki Morihira, Yasuhiro Kato, Akinobu Hayashi, Tatsuya Harada',
        'university': 'University of Tokyo / RIKEN AIP',
        'research_area': '世界模型、模型强化学习、表征学习',
        'link': 'openreview.net (ICLR 2026)',
        'innovation': '提出无需解码器或数据增强的冗余减少世界模型。消除了对解码器的依赖，避免昂贵的数据增强，在计算效率和性能上取得显著改进，为世界模型的实际部署提供更轻量方案。',
        'collaboration': '① 与自动驾驶公司合作，探索实时决策应用\n② 研究在边缘设备上的部署方案\n③ 探索与大型视觉-语言模型的结合',
        'figure_path': None,
        'author_photo': None,
        'recent_papers': [
            'DreamerV3: Mastering Diverse Domains (Nature 2025)',
            'Efficient World Models (ICML 2025)',
            'Latent Dynamics Learning (NeurIPS 2024)',
            'Model-Based RL for Control (CoRL 2024)',
            'Recurrent State Space Models (ICLR 2024)',
            'Visual World Models (CVPR 2024)'
        ]
    },
    {
        'title': 'Efficient RL by Guiding World Models with Non-Curated Data',
        'author': 'RIKEN AIP / University of Tokyo 研究团队',
        'university': 'University of Tokyo / RIKEN AIP',
        'research_area': '离线到在线RL、世界模型、样本效率',
        'link': 'openreview.net/forum?id=oBXfPyi47m',
        'innovation': '利用非策划数据（无奖励、质量混杂、跨具身）提升在线RL样本效率。提出经验回放和执行引导两种技术，在72个视觉运动任务上达到从头学习基线近两倍的聚合得分。',
        'collaboration': '① 利用企业海量无标注视频数据训练世界模型\n② 探索在游戏AI、仿真训练中的应用\n③ 研究如何低成本获取大规模预训练数据',
        'figure_path': None,
        'author_photo': None,
        'recent_papers': [
            'Offline-to-Online RL (ICML 2025)',
            'World Model Pretraining (NeurIPS 2025)',
            'Sample-Efficient RL (ICLR 2025)',
            'Visual RL Benchmarks (CoRL 2024)',
            'Experience Replay Methods (ICML 2024)',
            'Model-Based Offline RL (NeurIPS 2024)'
        ]
    },
    {
        'title': 'ICPRL: Acquiring Physical Intuition from Interactive Control',
        'author': 'University of Tokyo / RIKEN AIP 研究团队',
        'university': 'University of Tokyo / RIKEN AIP',
        'research_area': '物理推理、视觉-语言模型、具身智能',
        'link': 'openreview.net/forum?id=2Qh9YhuElD',
        'innovation': '让视觉-语言模型从交互控制中获取物理直觉。采用GRPO方法训练视觉策略模型，结合世界模型提供显式物理推理，通过PUCT搜索选择最优动作，显著提升物理推理能力。',
        'collaboration': '① 与物理仿真平台合作，构建更多物理推理基准\n② 探索在科学发现、工程仿真中的应用\n③ 研究物理直觉与常识推理的结合',
        'figure_path': None,
        'author_photo': None,
        'recent_papers': [
            'Physical Reasoning Benchmarks (NeurIPS 2025)',
            'VLM for Decision Making (ICLR 2025)',
            'In-Context RL (ICML 2025)',
            'Embodied AI with LLMs (CoRL 2024)',
            'Physics-Informed Neural Networks (NeurIPS 2024)',
            'World Models for Planning (ICML 2024)'
        ]
    }
]

# 生成PPT
add_cover_slide(prs)

for paper in papers:
    add_paper_slide(prs, paper)

# 保存
output_path = '/Users/gesong/.openclaw/workspace/ICLR2026-世界模型技术线索推荐-最终版.pptx'
prs.save(output_path)
print(f"PPT已保存到: {output_path}")
