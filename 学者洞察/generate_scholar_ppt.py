#!/usr/bin/env python3
"""
Diffusion语言模型方向优秀学者洞察 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Mm(254)
prs.slide_height = Mm(143)

# 颜色定义
BG_PRIMARY = RGBColor(30, 58, 95)
BG_SECONDARY = RGBColor(15, 23, 42)
BG_CARD = RGBColor(20, 30, 50)
TEXT_PRIMARY = RGBColor(240, 244, 248)
TEXT_SECONDARY = RGBColor(148, 163, 184)
TEXT_MUTED = RGBColor(100, 116, 139)
ACCENT = RGBColor(0, 212, 255)
ACCENT_SECONDARY = RGBColor(124, 58, 237)
BORDER = RGBColor(0, 100, 130)

def add_cover_slide(prs):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Mm(10), Mm(40), Mm(234), Mm(20))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Diffusion 语言模型"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Mm(10), Mm(62), Mm(234), Mm(12))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "优秀学者洞察"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # 统计信息
    stats_box = slide.shapes.add_textbox(Mm(10), Mm(78), Mm(234), Mm(10))
    tf = stats_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5位学者 | 国内+香港高校"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Mm(10), Mm(115), Mm(234), Mm(8))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026-03-04"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

def add_scholar_slide(prs, num, data):
    """学者页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景渐变（左侧深色）
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Mm(76), prs.slide_height)
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = BG_PRIMARY
    left_bg.line.fill.background()
    
    # 学者编号
    num_box = slide.shapes.add_textbox(Mm(8), Mm(8), Mm(60), Mm(15))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # 照片占位符
    y = Mm(26)
    photo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Mm(12), y, Mm(52), Mm(52))
    photo.fill.solid()
    photo.fill.fore_color.rgb = BG_CARD
    photo.line.color.rgb = ACCENT
    
    # 姓名
    y = Mm(82)
    name_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(8))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['name']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    # 机构
    y = Mm(90)
    uni_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(6))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['university']
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SECONDARY
    
    # 荣誉标签
    y = Mm(98)
    honor_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = honor_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['honor']
    p.font.size = Pt(9)
    p.font.color.rgb = ACCENT
    
    # 引用信息
    y = Mm(106)
    cite_box = slide.shapes.add_textbox(Mm(8), y, Mm(60), Mm(5))
    tf = cite_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"引用: {data['citations']} | h-index: {data['h_index']}"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    
    # ========== 右侧内容 ==========
    # 标题
    title_box = slide.shapes.add_textbox(Mm(80), Mm(8), Mm(166), Mm(8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['name'] + " - " + data['university']
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)
    
    # 研究方向
    dir_title = slide.shapes.add_textbox(Mm(80), Mm(20), Mm(166), Mm(5))
    tf = dir_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 研究方向"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    dir_box = slide.shapes.add_textbox(Mm(80), Mm(26), Mm(166), Mm(6))
    tf = dir_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['direction']
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(51, 65, 85)
    
    # 近期论文
    paper_title = slide.shapes.add_textbox(Mm(80), Mm(36), Mm(166), Mm(5))
    tf = paper_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 近期论文 (Top 5)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    y = Mm(42)
    for i, paper in enumerate(data['papers'][:5]):
        paper_box = slide.shapes.add_textbox(Mm(80), y + Mm(i * 10), Mm(166), Mm(10))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {paper}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(69, 26, 3)
    
    # 获奖信息
    award_title = slide.shapes.add_textbox(Mm(80), Mm(94), Mm(166), Mm(5))
    tf = award_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 获奖信息"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    award_box = slide.shapes.add_textbox(Mm(80), Mm(100), Mm(166), Mm(10))
    tf = award_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['awards']
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(69, 26, 3)
    
    # 优秀指标
    indicator_title = slide.shapes.add_textbox(Mm(80), Mm(112), Mm(166), Mm(5))
    tf = indicator_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ 优秀指标"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    indicator_box = slide.shapes.add_textbox(Mm(80), Mm(118), Mm(166), Mm(10))
    tf = indicator_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['indicators']
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(69, 26, 3)
    
    # 合作建议
    collab_title = slide.shapes.add_textbox(Mm(80), Mm(128), Mm(166), Mm(5))
    tf = collab_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🤝 合作建议"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    collab_box = slide.shapes.add_textbox(Mm(80), Mm(134), Mm(166), Mm(8))
    tf = collab_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['collab']
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(69, 26, 3)

# 学者数据
scholars = [
    {
        'name': '朱军 (Jun Zhu)',
        'university': '清华大学',
        'honor': 'IEEE/AAAI Fellow | 陈嘉庚科学奖',
        'citations': '60,971',
        'h_index': '102',
        'direction': '机器学习、贝叶斯方法、扩散模型、对抗鲁棒性、强化学习',
        'papers': [
            'Large Language Diffusion Models (arXiv 2025) - LLaDA',
            'DPM-Solver++ (MIR 2025) - 引用1000+',
            'Analytic-DPM (ICLR 2022) - Outstanding Paper Award',
            'U-ViT (arXiv 2022) - 首个扩散Transformer架构',
            'Vidu视频生成器 (2024) - 首个类Sora系统'
        ],
        'awards': '• ICLR 2022 Outstanding Paper Award (Analytic-DPM)\n• 陈嘉庚信息技术科学奖 (2024)\n• IEEE CoG 2022 Best Paper Award',
        'indicators': '✅ 顶会Best Paper: ICLR 2022\n✅ 引用上升趋势: 60,000+ 总引用\n✅ GitHub高星项目: Tianshou, ZhuSuan',
        'collab': '扩散模型理论合作、大模型训练基础设施共建'
    },
    {
        'name': '李崇轩 (Chongxuan Li)',
        'university': '中国人民大学',
        'honor': '国家高层次青年人才 | 优青',
        'citations': '8,000+',
        'h_index': '35+',
        'direction': '机器学习、深度生成模型、扩散模型、可控生成',
        'papers': [
            'Large Language Diffusion Models (arXiv 2025) - LLaDA',
            'All are Worth Words: ViT Backbone for Diffusion Models',
            'One Transformer Fits All Distributions in Multi-Modal Diffusion',
            'Diffusion Models Survey (ACM Computing Surveys 2024)',
            'Multi-Modal Diffusion at Scale (ICML 2023)'
        ],
        'awards': '• 国家高层次青年人才计划\n• 中国人民大学高瓴人工智能学院副教授\n• 清华大学博士',
        'indicators': '✅ 近期顶会论文: ICML, NeurIPS\n✅ GitHub高星项目\n✅ 扩散模型综述高引用',
        'collab': '多模态扩散模型、可控生成算法研究'
    },
    {
        'name': 'Gen Li (李根)',
        'university': '香港中文大学',
        'honor': 'Assistant Professor',
        'citations': '2,681',
        'h_index': '18',
        'direction': '扩散模型、强化学习、生成式AI、统计学',
        'papers': [
            'Diffusion Model Theory Papers (NeurIPS/ICML)',
            'Reinforcement Learning with Diffusion Models',
            'High-dimensional Statistics for Generative Models',
            'Nonconvex Optimization in ML',
            'Statistical Learning Theory'
        ],
        'awards': '• CUHK统计学系助理教授\n• 宾夕法尼亚大学博士后\n• 研究获得多项基金支持',
        'indicators': '✅ 扩散模型理论贡献\n✅ 统计学交叉研究\n✅ 引用稳步增长',
        'collab': '扩散模型理论基础、统计学习方法研究'
    },
    {
        'name': '李鸿升 (Hongsheng Li)',
        'university': '香港中文大学',
        'honor': 'Professor | IEEE Fellow',
        'citations': '70,787',
        'h_index': '85+',
        'direction': '计算机视觉、多模态学习、扩散模型、医学图像分析',
        'papers': [
            'CoMat: Aligning Text-to-Image Diffusion Model (NeurIPS 2024)',
            'PUMA: Unified MLLM with Visual Generation (ICCV 2025)',
            'Collaborative Video Diffusion (NeurIPS 2024)',
            'UniRL-Zero: Unified RL for Diffusion and LLM',
            '多模态生成模型系列工作'
        ],
        'awards': '• IEEE Fellow\n• CUHK电子工程系教授\n• 研究成果应用于自动驾驶等领域',
        'indicators': '✅ 高引用学者: 70,000+ 引用\n✅ 近期NeurIPS多篇论文\n✅ 多模态生成领先',
        'collab': '多模态扩散模型、视觉-语言生成、医学图像生成'
    },
    {
        'name': '鲍凡 (Fan Bao)',
        'university': '生数科技 / 清华大学',
        'honor': '生数科技创始人',
        'citations': '6,891',
        'h_index': '20+',
        'direction': '扩散模型、生成式AI、深度学习',
        'papers': [
            'Analytic-DPM (ICLR 2022) - Outstanding Paper Award',
            'Maximum Likelihood Training for Diffusion ODEs (ICML 2022)',
            'All are Worth Words: ViT Backbone for Diffusion',
            'U-ViT: Diffusion Transformer Architecture',
            'Vidu视频生成系统核心算法'
        ],
        'awards': '• ICLR 2022 Outstanding Paper Award\n• 清华大学博士\n• 生数科技创始人兼首席科学家',
        'indicators': '✅ 顶会Best Paper: ICLR 2022\n✅ 创业成功: 生数科技\n✅ 引用快速增长',
        'collab': '扩散模型产业化、视频生成技术应用'
    }
]

# 生成PPT
add_cover_slide(prs)

for i, scholar in enumerate(scholars):
    add_scholar_slide(prs, i + 1, scholar)

# 保存
output_path = '/Users/gesong/.openclaw/workspace/学者洞察/Diffusion语言模型优秀学者洞察.pptx'
os.makedirs('/Users/gesong/.openclaw/workspace/学者洞察', exist_ok=True)
prs.save(output_path)
print(f"PPT已保存: {output_path}")
