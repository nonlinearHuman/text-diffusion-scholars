#!/usr/bin/env python3
"""
将 frontend-slides HTML 演示文稿转换为 PPT
保持相似的视觉风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义 - Neon Cyber Theme
BG_PRIMARY = RGBColor(10, 14, 23)
BG_SECONDARY = RGBColor(17, 24, 39)
BG_CARD = RGBColor(20, 30, 50)
TEXT_PRIMARY = RGBColor(240, 244, 248)
TEXT_SECONDARY = RGBColor(148, 163, 184)
TEXT_MUTED = RGBColor(100, 116, 139)
ACCENT = RGBColor(0, 212, 255)
ACCENT_SECONDARY = RGBColor(124, 58, 237)
BORDER = RGBColor(0, 100, 130)

def add_cover_slide(prs):
    """添加封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    
    # 发光效果1
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-1), Inches(5), Inches(5))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = ACCENT
    glow1.fill.fore_color.brightness = 0.7
    glow1.line.fill.background()
    
    # 发光效果2
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0), Inches(4), Inches(4), Inches(4))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ACCENT_SECONDARY
    glow2.fill.fore_color.brightness = 0.7
    glow2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ICLR 2026"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术线索推荐"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # 方向
    direction_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6))
    tf = direction_box.text_frame
    p = tf.paragraphs[0]
    p.text = "方向：世界模型 (World Models)"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_SECONDARY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "面向大模型技术专家与领导层"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

def add_paper_slide(prs, num, title, subtitle, author, uni, country_flag, research, link, papers, innovation, collab):
    """添加论文页面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    
    # 论文编号
    num_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(1.5), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    # 论文标题
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.35), Inches(10.5), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = "Arial"
    
    # 论文副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.95), Inches(10.5), Inches(0.4))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_SECONDARY
    p.font.name = "Arial"
    
    # 左侧作者卡片背景
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.5), Inches(4.0), Inches(5.6))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = BG_CARD
    left_bg.line.color.rgb = BORDER
    left_bg.adjustments[0] = 0.05
    
    # 头像圆形
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.7), Inches(0.9), Inches(0.9))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = ACCENT
    avatar.line.fill.background()
    
    # 国旗emoji
    flag_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.85), Inches(0.6), Inches(0.6))
    tf = flag_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = country_flag
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # 作者名
    author_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.75), Inches(2.7), Inches(0.35))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = author
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = "Arial"
    
    # 机构
    uni_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(2.7), Inches(0.3))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = uni
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Arial"
    
    # 研究方向
    y = 2.6
    research_title = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3.6), Inches(0.25))
    tf = research_title.text_frame
    p = tf.paragraphs[0]
    p.text = "研究方向"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    research_content = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.28), Inches(3.6), Inches(0.4))
    tf = research_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = research
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SECONDARY
    p.font.name = "Arial"
    
    # 论文链接
    y = 3.4
    link_title = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3.6), Inches(0.25))
    tf = link_title.text_frame
    p = tf.paragraphs[0]
    p.text = "论文链接"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    link_content = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.28), Inches(3.6), Inches(0.3))
    tf = link_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = link
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    # 近期代表作
    y = 4.2
    papers_title = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3.6), Inches(0.25))
    tf = papers_title.text_frame
    p = tf.paragraphs[0]
    p.text = "近期代表作"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    y += 0.28
    for i, paper in enumerate(papers[:4]):
        paper_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + i * 0.32), Inches(3.6), Inches(0.32))
        tf = paper_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"→ {paper}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_SECONDARY
        p.font.name = "Arial"
    
    # 右侧创新点卡片
    right_bg1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.5), Inches(8.4), Inches(2.4))
    right_bg1.fill.solid()
    right_bg1.fill.fore_color.rgb = BG_CARD
    right_bg1.line.color.rgb = BORDER
    right_bg1.adjustments[0] = 0.05
    
    # 创新点标题
    innov_title = slide.shapes.add_textbox(Inches(4.7), Inches(1.65), Inches(8.0), Inches(0.35))
    tf = innov_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 核心创新点"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    # 创新点内容
    innov_content = slide.shapes.add_textbox(Inches(4.7), Inches(2.05), Inches(7.9), Inches(1.7))
    tf = innov_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = innovation
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_SECONDARY
    p.font.name = "Arial"
    p.line_spacing = 1.3
    
    # 合作建议卡片
    right_bg2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(4.1), Inches(8.4), Inches(3.0))
    right_bg2.fill.solid()
    right_bg2.fill.fore_color.rgb = RGBColor(15, 30, 45)
    right_bg2.line.color.rgb = BORDER
    right_bg2.adjustments[0] = 0.05
    
    # 合作建议标题
    collab_title = slide.shapes.add_textbox(Inches(4.7), Inches(4.25), Inches(8.0), Inches(0.35))
    tf = collab_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🤝 合作交流建议"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    
    # 合作建议内容
    y = 4.65
    for i, item in enumerate(collab):
        collab_box = slide.shapes.add_textbox(Inches(4.7), Inches(y + i * 0.45), Inches(7.9), Inches(0.45))
        tf = collab_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"● {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_SECONDARY
        p.font.name = "Arial"

# 论文数据
papers = [
    {
        'title': 'Test-Time Mixture of World Models',
        'subtitle': 'TMoW: 动态环境下的具身智能自适应框架',
        'author': 'Honguk Woo 等',
        'uni': '成均馆大学 (SKKU)',
        'country': '🇰🇷',
        'research': '具身智能、世界模型、自适应学习',
        'link': 'arxiv.org/abs/2601.22647',
        'papers': [
            'Incremental Learning of Skills (NeurIPS 2024)',
            'Embodied CoT Distillation (2024)',
            'Global DAG Task Scheduler (2021)'
        ],
        'innovation': '提出TMoW框架，将MoE扩展到具身智能体。支持测试时动态更新路由函数，实现多粒度原型路由、测试时精炼和蒸馏混合增强。在VirtualHome、ALFWorld、RLBench上验证了强大的零样本适应能力，比SayCanPay提升27.21%。',
        'collab': [
            '动态环境下的具身智能应用合作（家庭机器人、仓储物流）',
            '与国内机器人企业合作，将TMoW应用于实际场景',
            '探索在大模型决策系统中的集成应用'
        ]
    },
    {
        'title': 'Cross-Embodiment Offline RL',
        'subtitle': '异构机器人数据集的跨具身离线强化学习',
        'author': 'Haruki Abe 等',
        'uni': '东京大学 / RIKEN AIP',
        'country': '🇯🇵',
        'research': '跨具身学习、离线强化学习、机器人控制',
        'link': 'arxiv.org/abs/2602.18025',
        'papers': [
            'Offline RL for Robotics (CoRL 2025)',
            'Multi-Embodiment Transfer (NeurIPS 2024)',
            'Visual RL (CVPR 2024)'
        ],
        'innovation': '解决跨具身离线RL中如何利用形态各异的机器人数据学习通用控制先验的问题。提出基于形态相似性的分组策略，减少机器人间的梯度冲突，从异构次优轨迹中提取通用控制知识。',
        'collab': [
            '与国内机器人厂商合作，利用真实产线数据训练',
            '探索在工业机器人、服务机器人中的应用',
            '研究跨具身迁移学习的理论边界'
        ]
    },
    {
        'title': 'R2-Dreamer',
        'subtitle': '无解码器的冗余减少世界模型',
        'author': 'Naoki Morihira 等',
        'uni': '东京大学 / RIKEN AIP',
        'country': '🇯🇵',
        'research': '世界模型、模型强化学习、表征学习',
        'link': 'openreview.net (ICLR 2026)',
        'papers': [
            'DreamerV3 (Nature 2025)',
            'Efficient World Models (ICML 2025)',
            'Latent Dynamics (NeurIPS 2024)'
        ],
        'innovation': '提出无需解码器或数据增强的冗余减少世界模型。消除了对解码器的依赖，避免昂贵的数据增强，在计算效率和性能上取得显著改进，为世界模型的实际部署提供更轻量方案。',
        'collab': [
            '与自动驾驶公司合作，探索实时决策应用',
            '研究在边缘设备上的部署方案',
            '探索与大型视觉-语言模型的结合'
        ]
    },
    {
        'title': 'Guiding World Models',
        'subtitle': '利用非策划数据提升在线RL样本效率',
        'author': 'RIKEN AIP 团队',
        'uni': '东京大学 / RIKEN AIP',
        'country': '🇯🇵',
        'research': '离线到在线RL、世界模型、样本效率',
        'link': 'openreview.net/forum?id=oBXfPyi47m',
        'papers': [
            'Offline-to-Online RL (ICML 2025)',
            'World Model Pretraining (NeurIPS 2025)',
            'Sample-Efficient RL (ICLR 2025)'
        ],
        'innovation': '利用非策划数据（无奖励、质量混杂、跨具身）提升在线RL样本效率。提出经验回放和执行引导两种技术，在72个视觉运动任务上达到从头学习基线近两倍的聚合得分。',
        'collab': [
            '利用企业海量无标注视频数据训练世界模型',
            '探索在游戏AI、仿真训练中的应用',
            '研究如何低成本获取大规模预训练数据'
        ]
    },
    {
        'title': 'ICPRL',
        'subtitle': '从交互控制中获取物理直觉',
        'author': 'RIKEN AIP 团队',
        'uni': '东京大学 / RIKEN AIP',
        'country': '🇯🇵',
        'research': '物理推理、视觉-语言模型、具身智能',
        'link': 'openreview.net/forum?id=2Qh9YhuElD',
        'papers': [
            'Physical Reasoning Benchmarks (NeurIPS 2025)',
            'VLM for Decision Making (ICLR 2025)',
            'In-Context RL (ICML 2025)'
        ],
        'innovation': '让视觉-语言模型从交互控制中获取物理直觉。采用GRPO方法训练视觉策略模型，结合世界模型提供显式物理推理，通过PUCT搜索选择最优动作，显著提升物理推理能力。',
        'collab': [
            '与物理仿真平台合作，构建更多物理推理基准',
            '探索在科学发现、工程仿真中的应用',
            '研究物理直觉与常识推理的结合'
        ]
    }
]

# 生成PPT
add_cover_slide(prs)

for i, paper in enumerate(papers):
    add_paper_slide(
        prs,
        i + 1,
        paper['title'],
        paper['subtitle'],
        paper['author'],
        paper['uni'],
        paper['country'],
        paper['research'],
        paper['link'],
        paper['papers'],
        paper['innovation'],
        paper['collab']
    )

# 保存
output_path = '/Users/gesong/.openclaw/workspace/ICLR2026-WorldModels/ICLR2026-世界模型技术线索推荐.pptx'
prs.save(output_path)
print(f"PPT已保存到: {output_path}")
